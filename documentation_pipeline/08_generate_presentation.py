import os
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("python-pptx is required. pip install python-pptx")
    exit(1)

ROOT = Path(__file__).resolve().parent.parent
DOCS_DIR = ROOT / "docs"
PPT_PATH = DOCS_DIR / "QualiVision_AI_Presentation.pptx"

def create_presentation():
    print("="*50)
    print("Phase 8: Generating PPTX Presentation")
    print("="*50)
    
    prs = Presentation()
    
    # 1. Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "QualiVision AI"
    subtitle.text = "Industrial Quality Control System\nDeveloped by Jiphin George\nMCA Internship Project"
    
    # 2. Problem Statement
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Problem Statement"
    tf = body_shape.text_frame
    tf.text = "Manual quality inspection is:"
    
    p = tf.add_paragraph()
    p.text = "Slow and bottlenecks production."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Inconsistent due to human visual fatigue."
    p.level = 1
    
    # 3. Proposed Solution
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Proposed Solution: QualiVision AI"
    tf = body_shape.text_frame
    tf.text = "An automated Deep Learning CV platform featuring:"
    
    p = tf.add_paragraph()
    p.text = "EfficientNetV2B0 model for 99.8% accurate classification."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Real-time WebSocket streaming from factory cameras."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Grad-CAM explainability for transparent AI decisions."
    p.level = 1
    
    prs.save(str(PPT_PATH))
    print(f"[OK] Generated {PPT_PATH.name}")

if __name__ == "__main__":
    create_presentation()
